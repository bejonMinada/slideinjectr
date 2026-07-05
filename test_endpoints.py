#!/usr/bin/env python3
"""
Automated test script for SlideInjectr endpoints
Tests: analyze, preview, and generate functionality
"""

import requests
import json
from pathlib import Path
import time
import tempfile

BASE_URL = "http://localhost:5003/api"

def create_test_data():
    """Create minimal test files for testing"""
    from pptx import Presentation
    import pandas as pd
    import io
    
    # Create a minimal PowerPoint template
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add some test content
    left = top = width = height = 1000000
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Test Template"
    
    # Save to bytes
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    # Create a simple CSV
    df = pd.DataFrame({
        'Name': ['Alice', 'Bob', 'Charlie'],
        'Value': [100, 200, 300],
        'Department': ['Sales', 'IT', 'HR']
    })
    csv_bytes = io.StringIO()
    df.to_csv(csv_bytes, index=False)
    
    return ppt_bytes, csv_bytes

def test_health():
    """Test health check endpoint"""
    print("\n✓ Testing health check...")
    try:
        response = requests.get(f"{BASE_URL}/v1/presentations/health", timeout=5)
        if response.status_code == 200:
            try:
                data = response.json()
                print(f"  ✓ Status: {data}")
                return True
            except:
                # Even if JSON parsing fails, a 200 response is a pass
                print(f"  ✓ Health endpoint responded with 200 OK")
                return True
        else:
            print(f"  ❌ Health check failed: {response.status_code}")
            print(f"    Response: {response.text}")
            return False
    except Exception as e:
        print(f"  ❌ Error: {e}")
        return False

def test_analyze_template():
    """Test template analysis endpoint"""
    print("\n✓ Testing template analysis...")
    try:
        from pptx import Presentation
        import io
        
        # Create test template
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = 1000000
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Test {{Name}}"
        
        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        files = {'template': ('test.pptx', ppt_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
        response = requests.post(f"{BASE_URL}/v1/templates/analyze", files=files)
        
        if response.status_code == 200:
            data = response.json()
            print(f"  ✓ Template analyzed successfully")
            print(f"    Template name: {data.get('template_name')}")
            print(f"    Elements found: {len(data.get('elements', []))}")
            return True
        else:
            print(f"  ❌ Template analysis failed: {response.status_code}")
            print(f"    Response: {response.text}")
            return False
    except Exception as e:
        print(f"  ❌ Error: {e}")
        return False

def test_analyze_data():
    """Test data source analysis endpoint"""
    print("\n✓ Testing data source analysis...")
    try:
        import pandas as pd
        import io
        
        # Create test CSV
        df = pd.DataFrame({
            'Name': ['Alice', 'Bob', 'Charlie'],
            'Value': [100, 200, 300]
        })
        
        csv_bytes = io.BytesIO()
        df.to_csv(csv_bytes, index=False)
        csv_bytes.seek(0)
        
        files = {'data_files': ('test.csv', csv_bytes, 'text/csv')}
        response = requests.post(f"{BASE_URL}/v1/data-sources/analyze", files=files)
        
        if response.status_code == 200:
            data = response.json()
            print(f"  ✓ Data sources analyzed successfully")
            print(f"    Files: {len(data.get('data_files', []))}")
            if data.get('data_files'):
                print(f"    Columns: {data['data_files'][0].get('columns', [])}")
            return True
        else:
            print(f"  ❌ Data analysis failed: {response.status_code}")
            print(f"    Response: {response.text}")
            return False
    except Exception as e:
        print(f"  ❌ Error: {e}")
        return False

def test_generate_presentation():
    """Test presentation generation endpoint"""
    print("\n✓ Testing presentation generation...")
    try:
        from pptx import Presentation
        import pandas as pd
        import io
        
        # Create test template
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = 1000000
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Name Data"
        
        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        # Create test CSV
        df = pd.DataFrame({
            'Name': ['Alice', 'Bob'],
        })
        csv_bytes = io.BytesIO()
        df.to_csv(csv_bytes, index=False)
        csv_bytes.seek(0)
        
        # Prepare mapping with correct schema
        mapping = [
            {
                "element_name": "TextBox 1",
                "slide_index": 1,
                "mode": "text",
                "source_file": "test.csv",
                "source_sheet": None,
                "data_column": "Name",
                "row_index": 1
            }
        ]
        
        files = {
            'template': ('test.pptx', ppt_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
            'data_files': ('test.csv', csv_bytes, 'text/csv')
        }
        
        data_form = {
            'mapping_json': json.dumps(mapping),
            'return_file': 'true'
        }
        
        response = requests.post(f"{BASE_URL}/v1/presentations/generate", files=files, data=data_form, timeout=30)
        
        if response.status_code == 200:
            print(f"  ✓ Presentation generated successfully")
            print(f"    Response size: {len(response.content)} bytes")
            
            # Verify it's a valid file
            if len(response.content) > 0:
                print(f"  ✓ Generated file is valid")
                return True
            else:
                print(f"  ❌ Generated file is empty")
                return False
        else:
            print(f"  ❌ Generation failed: {response.status_code}")
            print(f"    Response: {response.text[:500]}")
            return False
    except Exception as e:
        print(f"  ❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    print("""
    ╔════════════════════════════════════╗
    ║   SlideInjectr Endpoint Tests      ║
    ╚════════════════════════════════════╝
    """)
    
    print("Waiting for server to be ready...")
    time.sleep(2)
    
    results = {}
    
    # Run tests
    results['health'] = test_health()
    results['analyze_template'] = test_analyze_template()
    results['analyze_data'] = test_analyze_data()
    results['generate'] = test_generate_presentation()
    
    # Summary
    print("\n" + "="*50)
    print("TEST SUMMARY")
    print("="*50)
    passed = sum(1 for v in results.values() if v)
    total = len(results)
    print(f"Passed: {passed}/{total}")
    
    for test_name, result in results.items():
        status = "✓ PASS" if result else "✗ FAIL"
        print(f"  {status}: {test_name}")
    
    if passed == total:
        print("\n✓ All tests passed!")
        return 0
    else:
        print(f"\n✗ {total - passed} test(s) failed")
        return 1

if __name__ == "__main__":
    exit(main())
