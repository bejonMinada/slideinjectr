from collections.abc import Iterable
from pathlib import Path

from pptx import Presentation

from pptx.oxml.ns import qn

from app.schemas.presentation import TemplateElement


def _shape_type_name(shape) -> str:
    try:
        return str(shape.shape_type).split(".")[-1]
    except Exception:
        return "UNKNOWN"


def _graphic_data_uri(shape) -> str | None:
    try:
        graphic_data = shape.element.graphic.graphicData
        return graphic_data.get("uri") if graphic_data is not None else None
    except Exception:
        return None


def _iter_shapes_recursively(shape_collection) -> Iterable:
    for shape in shape_collection:
        yield shape
        # Group shapes expose nested shapes through a 'shapes' collection.
        if hasattr(shape, "shapes") and shape.shapes is not None:
            yield from _iter_shapes_recursively(shape.shapes)


def analyze_template(template_path: Path) -> list[TemplateElement]:
    prs = Presentation(str(template_path))
    elements: list[TemplateElement] = []

    for slide_index, slide in enumerate(prs.slides):
        for shape in _iter_shapes_recursively(slide.shapes):
            placeholder_text = None
            has_text_frame = bool(getattr(shape, "has_text_frame", False))
            has_chart = bool(getattr(shape, "has_chart", False))
            has_table = bool(getattr(shape, "has_table", False))
            has_smartart = _graphic_data_uri(shape) == "http://schemas.openxmlformats.org/drawingml/2006/diagram"
            is_placeholder = bool(getattr(shape, "is_placeholder", False))
            is_mappable = has_text_frame or has_chart or has_table or has_smartart or is_placeholder

            if has_text_frame and getattr(shape, "text", None):
                placeholder_text = shape.text.strip()[:120] or None

            elements.append(
                TemplateElement(
                    slide_index=slide_index,
                    shape_id=getattr(shape, "shape_id", -1),
                    element_name=getattr(shape, "name", f"shape_{slide_index}"),
                    element_type=_shape_type_name(shape),
                    has_text_frame=has_text_frame,
                    has_chart=has_chart,
                    has_table=has_table,
                    has_smartart=has_smartart,
                    is_placeholder=is_placeholder,
                    is_mappable=is_mappable,
                    placeholder_text=placeholder_text,
                )
            )

    return elements
