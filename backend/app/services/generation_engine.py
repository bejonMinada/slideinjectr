from pathlib import Path
from uuid import uuid4

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData

from app.schemas.presentation import MappingRule
from app.services.data_source_analyzer import CSV_SHEET_NAME


def _read_data_file(path: Path) -> dict[str, pd.DataFrame]:
    suffix = path.suffix.lower()
    if suffix == ".csv":
        return {CSV_SHEET_NAME: pd.read_csv(path)}

    excel_sheets = pd.read_excel(path, sheet_name=None)
    return {sheet_name: df for sheet_name, df in excel_sheets.items()}


def _load_data_sources(data_sources: dict[str, Path]) -> dict[str, dict[str, pd.DataFrame]]:
    return {file_name: _read_data_file(path) for file_name, path in data_sources.items()}


def _pick_dataframe(
    loaded_sources: dict[str, dict[str, pd.DataFrame]],
    rule: MappingRule,
) -> pd.DataFrame | None:
    if not rule.source_file:
        return None
    if rule.source_file not in loaded_sources:
        return None

    file_sheets = loaded_sources[rule.source_file]
    if not file_sheets:
        return None

    if rule.source_sheet and rule.source_sheet in file_sheets:
        return file_sheets[rule.source_sheet]

    return next(iter(file_sheets.values()))


def _apply_table_replacement(shape, df: pd.DataFrame) -> None:
    if not getattr(shape, "has_table", False):
        return

    table = shape.table
    col_count = len(table.columns)
    row_count = len(table.rows)

    # First row gets headers, following rows get data.
    header_values = [str(col) for col in df.columns[:col_count]]
    for col_idx, header in enumerate(header_values):
        table.cell(0, col_idx).text = header

    data_rows = min(max(row_count - 1, 0), len(df.index))
    for out_row in range(data_rows):
        source_row = df.iloc[out_row]
        for col_idx in range(min(col_count, len(df.columns))):
            value = source_row.iloc[col_idx]
            table.cell(out_row + 1, col_idx).text = "" if pd.isna(value) else str(value)


def _apply_chart_replacement(shape, df: pd.DataFrame, rule: MappingRule) -> None:
    if not getattr(shape, "has_chart", False):
        return

    if rule.x_column and rule.x_column in df.columns:
        category_col = rule.x_column
    else:
        if len(df.columns) < 2:
            return
        category_col = df.columns[0]

    if rule.y_columns:
        value_cols = [column for column in rule.y_columns if column in df.columns and column != category_col]
    else:
        value_cols = [column for column in df.columns if column != category_col]

    if not value_cols:
        return

    chart_data = CategoryChartData()
    chart_data.categories = [str(value) for value in df[category_col].fillna("").tolist()]

    for column in value_cols:
        numeric_series = pd.to_numeric(df[column], errors="coerce").fillna(0)
        chart_data.add_series(str(column), numeric_series.tolist())

    try:
        shape.chart.replace_data(chart_data)
    except Exception:
        # Skip unsupported chart types to avoid breaking generation.
        return


def _apply_mappings(
    prs: Presentation,
    rules: list[MappingRule],
    loaded_sources: dict[str, dict[str, pd.DataFrame]],
) -> None:

    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for rule in rules:
                if rule.slide_index != slide_index:
                    continue
                if shape.name != rule.element_name:
                    continue
                if rule.mode == "keep":
                    continue

                if rule.mode == "manual":
                    if getattr(shape, "has_text_frame", False):
                        shape.text = rule.manual_text or ""
                    continue

                df = _pick_dataframe(loaded_sources, rule)
                if df is None or df.empty:
                    continue

                if rule.mode == "table" and getattr(shape, "has_table", False):
                    _apply_table_replacement(shape, df)
                    continue

                if rule.mode == "chart" and getattr(shape, "has_chart", False):
                    _apply_chart_replacement(shape, df, rule)
                    continue

                if rule.mode != "text" or not getattr(shape, "has_text_frame", False):
                    continue
                if not rule.data_column or rule.data_column not in df.columns:
                    continue
                zero_based_row = rule.row_index - 1
                if zero_based_row < 0 or zero_based_row >= len(df.index):
                    continue

                value = df.iloc[zero_based_row][rule.data_column]
                shape.text = "" if pd.isna(value) else str(value)


def apply_mapping_to_template(
    template_path: Path,
    data_sources: dict[str, Path],
    mapping_rules: list[MappingRule],
    output_dir: Path,
) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    loaded_sources = _load_data_sources(data_sources)

    prs = Presentation(str(template_path))

    _apply_mappings(prs, mapping_rules, loaded_sources)

    output_path = output_dir / f"generated-{uuid4().hex}.pptx"
    prs.save(str(output_path))
    return output_path
